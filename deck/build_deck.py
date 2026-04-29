r"""Build MedCite_Pitch.pptx from PITCH.md content.

Usage (from repo root):
    .\backend\.venv\Scripts\python.exe deck\build_deck.py

Slide content is locked to PITCH.md sections (see RESUME_PROMPT.txt).
Screenshots auto-load from deck/screenshots/ if present, else show
labeled placeholders. Drop the PNGs and re-run to refresh.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ----- Brand palette (sourced from frontend/public/icon.svg + Tailwind) -----
ACCENT = RGBColor(0x02, 0x84, 0xC7)      # sky-600 — MedCite brand
ACCENT_DARK = RGBColor(0x07, 0x59, 0x85)  # sky-700ish for hover/contrast
INK = RGBColor(0x0F, 0x17, 0x2A)          # slate-900 — primary text
MUTED = RGBColor(0x64, 0x74, 0x8B)        # slate-500 — secondary text
RULE = RGBColor(0xCB, 0xD5, 0xE1)         # slate-300 — rules / borders
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)      # slate-50 — card backgrounds
GREEN = RGBColor(0x05, 0x96, 0x69)        # emerald-600 — "verified"
AMBER = RGBColor(0xD9, 0x77, 0x06)        # amber-600 — "abstain"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DECK_DIR = Path(__file__).parent
SHOTS_DIR = DECK_DIR / "screenshots"
OUT_PATH = DECK_DIR / "MedCite_Pitch.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --------------------------- low-level helpers ---------------------------

def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank


def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Calibri",
    line_spacing=1.15,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, size=20, color=INK, bullet="•", line_spacing=1.4):
    """items: list of strings OR list of (head, tail) tuples for bold-lead bullets."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        bullet_run = p.add_run()
        bullet_run.text = f"{bullet}  "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = ACCENT
        bullet_run.font.bold = True
        if isinstance(item, tuple):
            head, tail = item
            head_run = p.add_run()
            head_run.text = head
            head_run.font.name = "Calibri"
            head_run.font.size = Pt(size)
            head_run.font.bold = True
            head_run.font.color.rgb = INK
            if tail:
                tail_run = p.add_run()
                tail_run.text = tail
                tail_run.font.name = "Calibri"
                tail_run.font.size = Pt(size)
                tail_run.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = item
            run.font.name = "Calibri"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_accent_bar(slide, x=Inches(0.6), y=Inches(0.6), w=Inches(0.18), h=Inches(0.55)):
    add_rect(slide, x, y, w, h, fill=ACCENT)


def add_footer(slide, label):
    add_rect(slide, Inches(0), Inches(7.18), SLIDE_W, Inches(0.04), fill=ACCENT)
    add_text(
        slide,
        Inches(0.6),
        Inches(7.22),
        Inches(8),
        Inches(0.3),
        "MedCite  ·  Jubilant Pharma Hackathon  ·  abhi04-medcite.vercel.app",
        size=10,
        color=MUTED,
    )
    add_text(
        slide,
        Inches(10.5),
        Inches(7.22),
        Inches(2.3),
        Inches(0.3),
        label,
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_screenshot_or_placeholder(slide, filename, x, y, w, h, label):
    path = SHOTS_DIR / filename
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        # subtle border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        border.shadow.inherit = False
        border.fill.background()
        border.line.color.rgb = RULE
        border.line.width = Pt(1)
        return
    # placeholder
    add_rect(slide, x, y, w, h, fill=SURFACE, line=RULE)
    add_text(
        slide,
        x,
        y + Inches(0.4),
        w,
        Inches(0.5),
        "[ screenshot ]",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        bold=True,
    )
    add_text(
        slide,
        x + Inches(0.3),
        y + Inches(1.0),
        w - Inches(0.6),
        Inches(1.4),
        f"Drop file at:\ndeck/screenshots/{filename}\n\n{label}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        line_spacing=1.4,
    )


def add_title(slide, title, kicker=None):
    add_accent_bar(slide)
    if kicker:
        add_text(
            slide,
            Inches(0.95),
            Inches(0.55),
            Inches(8),
            Inches(0.32),
            kicker.upper(),
            size=11,
            color=ACCENT,
            bold=True,
        )
        add_text(
            slide,
            Inches(0.95),
            Inches(0.85),
            Inches(11.5),
            Inches(0.7),
            title,
            size=32,
            color=INK,
            bold=True,
        )
    else:
        add_text(
            slide,
            Inches(0.95),
            Inches(0.6),
            Inches(11.5),
            Inches(0.8),
            title,
            size=34,
            color=INK,
            bold=True,
        )


# ------------------------------- slides --------------------------------

def slide_1_title(prs):
    s = add_blank_slide(prs)
    # left accent panel
    add_rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, fill=ACCENT)

    # logo mark — a simple "Rx in a rounded square" using brand color
    mark = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(2.0), Inches(0.95), Inches(0.95)
    )
    mark.shadow.inherit = False
    mark.fill.solid()
    mark.fill.fore_color.rgb = ACCENT
    mark.line.fill.background()
    add_text(
        s,
        Inches(0.95),
        Inches(2.05),
        Inches(0.95),
        Inches(0.85),
        "M",
        size=44,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        font="Calibri",
    )

    add_text(
        s, Inches(2.1), Inches(2.0), Inches(10), Inches(0.45),
        "MEDCITE", size=14, color=ACCENT, bold=True,
    )
    add_text(
        s, Inches(2.1), Inches(2.4), Inches(11), Inches(1.2),
        "Clinical Q&A that\nnever hallucinates.",
        size=54, color=INK, bold=True, line_spacing=1.05,
    )
    add_text(
        s, Inches(2.1), Inches(4.6), Inches(10.5), Inches(0.6),
        "Cited. Verified. Honest.",
        size=22, color=MUTED, italic=True,
    )

    # bottom meta strip
    add_rect(s, Inches(2.1), Inches(5.6), Inches(10.5), Inches(0.04), fill=RULE)
    add_text(
        s, Inches(2.1), Inches(5.75), Inches(10.5), Inches(0.4),
        "Jubilant Pharma Hackathon  ·  3-day build  ·  1 person",
        size=14, color=MUTED,
    )
    add_text(
        s, Inches(2.1), Inches(6.15), Inches(10.5), Inches(0.4),
        "https://abhi04-medcite.vercel.app",
        size=14, color=ACCENT, bold=True,
    )


def slide_2_problem(prs):
    s = add_blank_slide(prs)
    add_title(s, "Doctors need answers they can trust.", kicker="The problem")

    add_bullets(
        s, Inches(0.95), Inches(1.95), Inches(11.5), Inches(4.5),
        [
            ("Time-starved.  ",
             "Doctors need fast, trustworthy medical answers at the point of care."),
            ("ChatGPT hallucinates citations.  ",
             "Fake DOIs, made-up authors, papers that don't exist. Documented and unsafe."),
            ("PubMed isn't a tool.  ",
             "200 raw results, no synthesis — doctors have to read everything themselves."),
            ("\"AI doctor\" tools are chatbots.  ",
             "Optimized for engagement. Doctors don't want a conversation."),
        ],
        size=20, line_spacing=1.5,
    )

    # gap callout box
    add_rect(s, Inches(0.95), Inches(5.7), Inches(11.5), Inches(1.15),
             fill=SURFACE, line=ACCENT)
    add_text(
        s, Inches(1.2), Inches(5.85), Inches(0.5), Inches(0.85),
        "→", size=32, color=ACCENT, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, Inches(1.85), Inches(5.85), Inches(10.4), Inches(0.4),
        "THE GAP",
        size=11, color=ACCENT, bold=True,
    )
    add_text(
        s, Inches(1.85), Inches(6.18), Inches(10.4), Inches(0.6),
        "No product combines LLM speed and synthesis with the auditability of a real literature search.",
        size=16, color=INK, bold=True,
    )

    add_footer(s, "2 / 10")


def slide_3_what_we_built(prs):
    s = add_blank_slide(prs)
    add_title(s, "One question. One cited answer card. Three seconds.",
              kicker="What we built")

    # Left column: features
    add_bullets(
        s, Inches(0.95), Inches(1.95), Inches(6.0), Inches(5.2),
        [
            ("2–4 sentence answer  ", "every claim tagged [1] [2] [3]"),
            ("5 source cards  ",
             "title, journal, year, authors, evidence-level badge"),
            ("Quoted passage  ",
             "the exact chunk the LLM read — verify at a glance"),
            ("Clickable PubMed + DOI  ",
             "real URLs, stitched from PMIDs (LLM never writes them)"),
            ("Confidence meter  ", "green ≥ 0.75, otherwise abstain"),
            ("\"Verified KB\" badge  ",
             "or \"Live multi-AI\" if escalated to PubMed"),
        ],
        size=16, line_spacing=1.5,
    )

    # Right column: screenshot
    add_screenshot_or_placeholder(
        s,
        "hero2-empagliflozin.png",
        Inches(7.2),
        Inches(1.95),
        Inches(5.5),
        Inches(4.6),
        label="Open https://abhi04-medcite.vercel.app, click hero #2 (empagliflozin / HFpEF), wait for the answer card with RCT + Review badges, screenshot the answer + first 2 source cards.",
    )

    # Caption under screenshot
    add_text(
        s, Inches(7.2), Inches(6.65), Inches(5.5), Inches(0.4),
        "Hero query #2 — empagliflozin in HFpEF  ·  Tier-1 cache hit  ·  conf 0.80",
        size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "3 / 10")


def slide_4_trust_pitch(prs):
    s = add_blank_slide(prs)
    add_title(s, "The whole pitch in four sentences.", kicker="The trust pitch")

    # Big quote card
    add_rect(s, Inches(0.95), Inches(1.95), Inches(11.5), Inches(4.6),
             fill=SURFACE, line=ACCENT)
    # huge opening quote mark
    add_text(
        s, Inches(1.2), Inches(1.85), Inches(0.8), Inches(1.0),
        "\u201C", size=80, color=ACCENT, bold=True,
    )

    quote_lines = [
        "Every other AI tool starts with the LLM and asks it to remember sources.",
        "We start with PubMed and ask the LLM to summarize what's actually there.",
        "The LLM is never allowed to invent a citation, never allowed to answer from",
        "memory, and a second model from a different company verifies every claim",
        "before the doctor ever sees it. If it can't reach 75% confidence — it says so.",
        "",
        "That's the difference between a chatbot and a clinical tool.",
    ]
    add_text(
        s, Inches(1.5), Inches(2.55), Inches(10.5), Inches(3.6),
        "\n".join(quote_lines),
        size=20, color=INK, line_spacing=1.45,
    )

    # attribution / "memorize this"
    add_text(
        s, Inches(0.95), Inches(6.7), Inches(11.5), Inches(0.4),
        "≈ 25 seconds spoken  ·  memorize verbatim",
        size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "4 / 10")


def slide_5_seven_rules(prs):
    s = add_blank_slide(prs)
    add_title(s, "What MedCite will never do.", kicker="The 7 hard rules")

    rules = [
        ("LLMs never write URLs.",
         "Backend stitches them from PMIDs in metadata."),
        ("LLMs never answer from memory.",
         "Chunks-only, or output INSUFFICIENT_EVIDENCE."),
        ("A different vendor verifies.",
         "Google Gemini synthesizes; Meta Llama verifies."),
        ("Confidence below 0.75 → abstain.",
         "No low-confidence guesses. Ever."),
        ("The doctor controls the live search.",
         "Cache runs auto. Live PubMed runs only on click."),
        ("Every citation has the quoted passage.",
         "Doctors verify the source at a glance."),
        ("It is a clinical tool, not a chatbot.",
         "One question, one answer card. No history."),
    ]

    # 2-column grid: 4 rows × 2 cols (last cell on row 4 is the closing tag)
    cols = 2
    col_w = Inches(5.65)
    row_h = Inches(1.15)
    gap_x = Inches(0.2)
    start_x = Inches(0.95)
    start_y = Inches(1.95)

    for i, (head, tail) in enumerate(rules):
        col = i % cols
        row = i // cols
        x = start_x + (col_w + gap_x) * col
        y = start_y + row_h * row
        # number badge
        badge = s.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y + Inches(0.18), Inches(0.55), Inches(0.55)
        )
        badge.shadow.inherit = False
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        add_text(
            s, x, y + Inches(0.18), Inches(0.55), Inches(0.55),
            str(i + 1), size=18, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # head + tail
        add_text(
            s, x + Inches(0.75), y + Inches(0.05), col_w - Inches(0.8), Inches(0.45),
            head, size=15, color=INK, bold=True,
        )
        add_text(
            s, x + Inches(0.75), y + Inches(0.5), col_w - Inches(0.8), Inches(0.55),
            tail, size=12, color=MUTED, line_spacing=1.3,
        )

    # closing pill in the empty 8th cell
    x = start_x + (col_w + gap_x)
    y = start_y + row_h * 3
    add_rect(s, x, y + Inches(0.1), col_w, Inches(0.85),
             fill=ACCENT)
    add_text(
        s, x, y + Inches(0.1), col_w, Inches(0.85),
        "These are the architecture — not marketing.",
        size=14, color=WHITE, bold=True, italic=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    add_footer(s, "5 / 10")


def slide_6_architecture(prs):
    s = add_blank_slide(prs)
    add_title(s, "Two-tier retrieval. Cross-vendor verification.",
              kicker="Architecture")

    # Top: doctor input
    cx = SLIDE_W / 2
    add_rect(s, cx - Inches(2.5), Inches(1.85), Inches(5), Inches(0.55),
             fill=INK)
    add_text(
        s, cx - Inches(2.5), Inches(1.85), Inches(5), Inches(0.55),
        "Doctor types a query",
        size=15, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Arrow down to retriever
    add_rect(s, cx - Emu(15000), Inches(2.45), Emu(30000), Inches(0.35),
             fill=MUTED)

    # Tier 1 box (full-width)
    add_rect(s, Inches(2.5), Inches(2.85), Inches(8.3), Inches(0.85),
             fill=ACCENT)
    add_text(
        s, Inches(2.5), Inches(2.9), Inches(8.3), Inches(0.4),
        "TIER 1 — LOCAL RETRIEVER  ·  ~3 seconds",
        size=11, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(2.5), Inches(3.25), Inches(8.3), Inches(0.4),
        "LanceDB  ·  17,456 chunks  ·  MiniLM-L6-v2 cosine search  ·  threshold 0.55",
        size=13, color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Branches: Found / Not found
    # Left branch (Found → synth+verify)
    add_text(
        s, Inches(2.0), Inches(3.85), Inches(3.5), Inches(0.4),
        "↓  Found  (sim ≥ 0.55)",
        size=12, color=INK, bold=True, align=PP_ALIGN.CENTER,
    )
    add_rect(s, Inches(1.7), Inches(4.3), Inches(4.1), Inches(2.2),
             fill=SURFACE, line=RULE)
    add_text(
        s, Inches(1.7), Inches(4.4), Inches(4.1), Inches(0.4),
        "SYNTHESIZE + VERIFY",
        size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER,
    )
    add_bullets(
        s, Inches(1.95), Inches(4.8), Inches(3.7), Inches(1.6),
        [
            ("Gemini 2.5 Flash-Lite  ", "synth, chunks-only"),
            ("Llama 3.3 70B (Groq)  ", "verify, JSON output"),
            ("conf ≥ 0.75  ", "→ answer, else abstain"),
        ],
        size=11, line_spacing=1.35,
    )

    # Right branch (Not found → live PubMed)
    add_text(
        s, Inches(7.8), Inches(3.85), Inches(3.5), Inches(0.4),
        "↓  Not found  (doctor clicks)",
        size=12, color=INK, bold=True, align=PP_ALIGN.CENTER,
    )
    add_rect(s, Inches(7.5), Inches(4.3), Inches(4.1), Inches(2.2),
             fill=SURFACE, line=GREEN)
    add_text(
        s, Inches(7.5), Inches(4.4), Inches(4.1), Inches(0.4),
        "TIER 2 — LIVE PUBMED  ·  ~20 sec",
        size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER,
    )
    add_bullets(
        s, Inches(7.75), Inches(4.8), Inches(3.7), Inches(1.6),
        [
            ("PubMed E-utilities  ", "esearch + efetch top-10"),
            ("Same synth + verify  ", "pipeline reused"),
            ("Write-back  ", "→ next ask hits Tier 1"),
        ],
        size=11, line_spacing=1.35,
    )

    add_footer(s, "6 / 10")


def slide_7_live_demo(prs):
    s = add_blank_slide(prs)
    add_title(s, "Live demo.", kicker="Switch to browser")

    # Left: query list + URL
    add_text(
        s, Inches(0.95), Inches(1.95), Inches(7.0), Inches(0.4),
        "FOUR HERO QUERIES",
        size=11, color=ACCENT, bold=True,
    )
    add_bullets(
        s, Inches(0.95), Inches(2.35), Inches(7.0), Inches(4.0),
        [
            ("#1 Empagliflozin / HFpEF  ",
             "Tier-1 hit, RCT + Review badges, conf 0.80"),
            ("#2 Drug-resistant TB 2024  ",
             "self-improvement story — yesterday missed, today instant"),
            ("#3 Levetiracetam status epilepticus  ",
             "live multi-AI: PubMed → Gemini → Llama, ~20 sec"),
            ("#4 Acetaminophen in 3rd trimester  ",
             "honest abstention — \"closest match scored 0.45\""),
        ],
        size=15, line_spacing=1.55,
    )

    # URL pill
    add_rect(s, Inches(0.95), Inches(6.35), Inches(7.0), Inches(0.65),
             fill=INK)
    add_text(
        s, Inches(0.95), Inches(6.35), Inches(7.0), Inches(0.65),
        "abhi04-medcite.vercel.app",
        size=18, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Right: abstention screenshot (the visual proof)
    add_text(
        s, Inches(8.2), Inches(1.95), Inches(4.5), Inches(0.4),
        "WHAT JUDGES WILL SEE ON #4",
        size=11, color=AMBER, bold=True,
    )
    add_screenshot_or_placeholder(
        s,
        "hero5-acetaminophen-abstain.png",
        Inches(8.2),
        Inches(2.35),
        Inches(4.5),
        Inches(4.2),
        label="Open hero #5 (acetaminophen / 3rd trimester pregnancy). Wait for the amber abstention screen with \"closest match scored 0.45\". Screenshot the full abstention card.",
    )
    add_text(
        s, Inches(8.2), Inches(6.65), Inches(4.5), Inches(0.4),
        "The abstention IS the feature.",
        size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "7 / 10")


def slide_8_numbers(prs):
    s = add_blank_slide(prs)
    add_title(s, "The numbers.", kicker="By the numbers")

    stats = [
        ("17,456", "chunks in the\npre-indexed cache"),
        ("~10,000", "PubMed abstracts\n(Diabetes + Cardiology)"),
        ("2", "LLM vendors for\ncross-verification"),
        ("0.75", "confidence threshold\nto surface an answer"),
        ("2–5 sec", "Tier-1 cache-hit\nresponse time"),
        ("15–25 sec", "live PubMed escalation\n(API-bound)"),
        ("0", "invented citations across\nall hero queries"),
        ("$0.001", "approximate LLM cost\nper query at scale"),
    ]

    # 4 × 2 grid
    cols = 4
    col_w = Inches(2.85)
    row_h = Inches(2.15)
    gap_x = Inches(0.18)
    gap_y = Inches(0.25)
    start_x = Inches(0.95)
    start_y = Inches(2.0)

    for i, (big, small) in enumerate(stats):
        col = i % cols
        row = i // cols
        x = start_x + (col_w + gap_x) * col
        y = start_y + (row_h + gap_y) * row
        add_rect(s, x, y, col_w, row_h, fill=SURFACE, line=RULE)
        # accent top stripe
        add_rect(s, x, y, col_w, Inches(0.08), fill=ACCENT)
        add_text(
            s, x, y + Inches(0.3), col_w, Inches(1.0),
            big, size=36, color=ACCENT, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x + Inches(0.15), y + Inches(1.3), col_w - Inches(0.3), Inches(0.8),
            small, size=11, color=INK,
            align=PP_ALIGN.CENTER, line_spacing=1.3,
        )

    # 3-day build tag bottom
    add_text(
        s, Inches(0.95), Inches(6.85), Inches(11.5), Inches(0.35),
        "3-day build  ·  1 person  ·  every successful escalation grows the cache",
        size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "8 / 12")


def slide_9_qa(prs):
    s = add_blank_slide(prs)
    add_title(s, "Top three questions we expect.", kicker="Q&A teaser")

    qas = [
        (
            "\"Why not just use GPT-4 with web search?\"",
            "GPT-4 cites sources to justify what it already wrote. We retrieve sources first, then ask the LLM to summarize only those sources. Our LLM is structurally incapable of citing something it didn't read.",
        ),
        (
            "\"What stops the LLM from making things up about the sources you give it?\"",
            "Two things. The synthesizer prompt forbids it — must output INSUFFICIENT_EVIDENCE if a chunk isn't relevant. And a different-vendor verifier (Llama 3.3 70B from Meta, via Groq) reads answer + sources and outputs {confidence, unsupported_claims}. Below 0.75 → abstain.",
        ),
        (
            "\"Is this just RAG?\"",
            "Yes — and \"just RAG\" is the correct architecture. The novelty is (a) the cross-vendor verifier as a hard safety gate, (b) the self-improvement write-back loop, and (c) the explicit abstention behavior. Most RAG demos still let the LLM speak when retrieval was weak. We don't.",
        ),
    ]

    y = Inches(1.95)
    for i, (q, a) in enumerate(qas):
        # Q badge
        add_rect(s, Inches(0.95), y, Inches(0.55), Inches(0.55), fill=ACCENT)
        add_text(
            s, Inches(0.95), y, Inches(0.55), Inches(0.55),
            "Q", size=20, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, Inches(1.7), y + Inches(0.02), Inches(11), Inches(0.45),
            q, size=16, color=INK, bold=True,
        )
        add_text(
            s, Inches(1.7), y + Inches(0.55), Inches(11), Inches(1.1),
            a, size=12, color=MUTED, line_spacing=1.45,
        )
        y += Inches(1.65)

    add_text(
        s, Inches(0.95), Inches(6.95), Inches(11.5), Inches(0.35),
        "Full Q&A bank: 12 rehearsed answers in PITCH.md §9",
        size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "9 / 12")


def slide_10_local_kb_benefits(prs):
    s = add_blank_slide(prs)
    add_title(s, "Why we start with a local knowledge base.",
              kicker="System design choice")

    add_bullets(
        s, Inches(0.95), Inches(1.95), Inches(11.5), Inches(4.9),
        [
            ("Fast and reliable by default.  ",
             "Most doctor questions are answered in seconds without waiting on live PubMed calls."),
            ("Not blocked by external API limits.  ",
             "Live PubMed can face rate limits, latency spikes, and dependency risk under load."),
            ("Scales better for repeated queries.  ",
             "Common questions hit local cache instantly instead of repeating external calls."),
            ("We control the data layer end-to-end.  ",
             "We decide what is indexed, how it is chunked, versioned, and quality-checked."),
            ("Security and governance are stronger.  ",
             "Data handling stays on controlled backend infrastructure with auditable retrieval rules."),
            ("Live search remains an explicit fallback.  ",
             "If local evidence is weak, doctor can trigger live PubMed plus verification."),
        ],
        size=15, line_spacing=1.45,
    )

    add_rect(s, Inches(0.95), Inches(6.55), Inches(11.5), Inches(0.42), fill=SURFACE, line=RULE)
    add_text(
        s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.3),
        "Tier-1 local-first architecture = speed, control, and resilience.",
        size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "10 / 12")


def slide_11_frontend_freedom(prs):
    s = add_blank_slide(prs)
    add_title(s, "One backend, many frontends.",
              kicker="Platform flexibility")

    add_bullets(
        s, Inches(0.95), Inches(1.95), Inches(11.5), Inches(4.9),
        [
            ("Backend is frontend-agnostic.  ",
             "FastAPI endpoints are independent of UI framework and can serve any client."),
            ("No backend rewrite for UI changes.  ",
             "We can redesign UX and interaction flow without changing retrieval or verifier logic."),
            ("Multi-platform ready by design.  ",
             "The same API powers our production web app and can power mobile or desktop clients."),
            ("Proof of portability.  ",
             "Alongside the web app, we also built a Flutter demo client (with current hackathon limits)."),
            ("Faster iteration with safer core logic.  ",
             "Frontend can move quickly while backend safety rules remain stable."),
            ("Future-ready integration layer.  ",
             "Hospital portals, internal tools, and partner apps can all use the same API core."),
        ],
        size=15, line_spacing=1.45,
    )

    add_rect(s, Inches(0.95), Inches(6.55), Inches(11.5), Inches(0.42), fill=SURFACE, line=RULE)
    add_text(
        s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.3),
        "One verified backend, many possible interfaces.",
        size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "11 / 12")


def slide_12_closing(prs):
    s = add_blank_slide(prs)
    # full-bleed accent panel left
    add_rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, fill=ACCENT)

    add_text(
        s, Inches(0.95), Inches(0.95), Inches(11), Inches(0.45),
        "CLOSING", size=12, color=ACCENT, bold=True,
    )

    closing_lines = [
        "This isn't a chatbot.",
        "It's a clinical tool.",
    ]
    add_text(
        s, Inches(0.95), Inches(1.45), Inches(11.8), Inches(2.2),
        "\n".join(closing_lines),
        size=60, color=INK, bold=True, line_spacing=1.05,
    )

    body = (
        "It cites every claim, it verifies with a second model from a "
        "different company, and when it doesn't know — it says so."
    )
    add_text(
        s, Inches(0.95), Inches(3.95), Inches(11.5), Inches(1.5),
        body, size=22, color=INK, line_spacing=1.4,
    )

    add_text(
        s, Inches(0.95), Inches(5.5), Inches(11.5), Inches(0.5),
        "That's the only safe way to put an LLM in front of a doctor.",
        size=20, color=ACCENT, bold=True, italic=True,
    )

    # CTA strip
    add_rect(s, Inches(0.95), Inches(6.3), Inches(11.5), Inches(0.7),
             fill=INK)
    add_text(
        s, Inches(0.95), Inches(6.3), Inches(11.5), Inches(0.7),
        "MedCite  ·  abhi04-medcite.vercel.app  ·  github.com/Abhishek0489/medCite",
        size=16, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    add_text(
        s, Inches(0.95), Inches(7.15), Inches(11.5), Inches(0.3),
        "Thank you. Questions?",
        size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )

    add_footer(s, "12 / 12")


# --------------------------------- main --------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_what_we_built(prs)
    slide_4_trust_pitch(prs)
    slide_5_seven_rules(prs)
    slide_6_architecture(prs)
    slide_7_live_demo(prs)
    slide_8_numbers(prs)
    slide_9_qa(prs)
    slide_10_local_kb_benefits(prs)
    slide_11_frontend_freedom(prs)
    slide_12_closing(prs)

    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({OUT_PATH.stat().st_size / 1024:.1f} KB)")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
